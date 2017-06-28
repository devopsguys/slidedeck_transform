import sys
import json
from pptx import Presentation
from pptx.util import Inches, Pt

def delete_slide(prs, slide):
    #Make dictionary with necessary information
    id_dict = {slide.id: [i, slide.rId] for i, slide in enumerate(prs.slides._sldIdLst)}
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

print 'Number of arguments:', len(sys.argv), 'arguments.'
print 'Argument List:', str(sys.argv)

PRS = Presentation(sys.argv[1])

width = height = Inches(1)
left = top = Inches(2)

for slide in PRS.slides:
    g = "/Users/edmundd/Desktop/logo.gif"
    top = PRS.slide_height - Inches(1)
    

    # print "==="
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    for shape in slide.placeholders:
        if shape.name == "Picture Placeholder 3":
            idx = shape.placeholder_format.idx
            pic = slide.shapes.add_picture(g, left, top, width, height)
            # print slide.placeholders[idx].element.xml
        # else:
            # print shape.name

    # print "==="

   

    if slide.has_notes_slide:
        text_frame = slide.notes_slide.notes_text_frame
        try:
            metadata = json.loads(text_frame.text)
            if "hello" in metadata['tags']:
                delete_slide(PRS, slide)
        except:
            pass

PRS.save('/Users/edmundd/Desktop/Test Presentation2.pptx')
