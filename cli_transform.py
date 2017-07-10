import os
import argparse
import json
from pptx import Presentation
from PIL import Image

import sdt_common
import sdt_tag_parse
import sdt_logo_stamp
import sdt_string_replace


def parse_args():
    parser = argparse.ArgumentParser(
        description='Transform a pptx presentation')
    required = parser.add_argument_group('Required arguments')
    optional = parser.add_argument_group('Optional arguments')
    optional.add_argument('--trim-logo', action='store_true',
                          help="Trim any whitespace from the logo")
    optional.add_argument('--out', action='store',
                          help="The path to the output file")
    required.add_argument('--file', nargs=1, action='store',
                          help="Path to the powerpoint presentation file")
    required.add_argument('--logo', nargs=1, action='store',
                          help="Path to the logo file")
    optional.add_argument('--tags', nargs="?", action='store', default="",
                          help="List of tags to remove")
    optional.add_argument('--templates', nargs=1, action='store',
                          help="JSON dictionary of template names and their values")

    return parser.parse_args()


def main():
    ARGS = parse_args()

    presentation_file = ARGS.file[0]
    presentation = Presentation(presentation_file)
    logo_image = ARGS.logo[0]
    temp_logo_file = sdt_common.TEMP_LOGO_FILE
    tags_to_delete = ARGS.tags.split(",")
    output_file = ARGS.out or presentation_file.replace(
        ".ppt", "-{0}.ppt".format("new"))

    all_tags = sdt_tag_parse.get_all_tags_in_presentation(presentation)
    print "All tags: " + str(all_tags)
    total_slides = len(presentation.slides)

    if ARGS.trim_logo:
        sdt_logo_stamp.trim(Image.open(logo_image)).save(temp_logo_file)
    else:
        Image.open(logo_image).save(temp_logo_file)

    for index, slide in enumerate(presentation.slides):
        print "{0}/{1}".format(index + 1, total_slides)

        sdt_logo_stamp.update_logo(slide)
        if(ARGS.templates):
            templates = json.loads(ARGS.templates[0])
            for template in templates:
                sdt_string_replace.replace_template_string(
                    slide, template['name'], template['value'])
        sdt_tag_parse.delete_slide_if_tag_matches(
            presentation, index, slide, tags_to_delete)

    os.remove(temp_logo_file)
    presentation.save(output_file)
    print "Done!"


if __name__ == '__main__':
    main()
