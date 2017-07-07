import os
import copy
import argparse
from pptx import Presentation
from PIL import Image

import sdt_common
import sdt_tag_parse
import sdt_logo_stamp
import sdt_string_replace


def parse_args():
    parser = argparse.ArgumentParser(
        description='Transform a pptx presentation')
    required = parser.add_argument_group('Required arguments')
    optional = parser.add_argument_group('Optional arguments')
    optional.add_argument('--trim-logo', action='store_true',
                          help="Trim any whitespace from the logo")
    required.add_argument('--file', nargs=1, action='store',
                          help="Path to the powerpoint presentation file")
    required.add_argument('--logo', nargs=1, action='store',
                          help="Path to the logo file")
    optional.add_argument('--tags', nargs="?", action='store', default="",
                          help="List of tags to remove")
    required.add_argument('--client', nargs=1,
                          action='store', help="The name of the client")

    return parser.parse_args()


def main():
    ARGS = parse_args()

    presentation_file = ARGS.file[0]
    presentation = Presentation(presentation_file)
    logo_image = ARGS.logo[0]
    temp_logo_file = sdt_common.TEMP_LOGO_FILE

    all_tags = sdt_tag_parse.get_all_tags_in_presentation(presentation)
    print "All tags: " + str(all_tags)
    total_slides = len(presentation.slides)

    if ARGS.trim_logo:
        sdt_logo_stamp.trim(Image.open(logo_image)).save(temp_logo_file)
    else:
        Image.open(logo_image).save(temp_logo_file)

    for index, slide in enumerate(presentation.slides):
        print "{0}/{1}".format(index + 1, total_slides)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                for template in sdt_string_replace.find_templates_in_string(paragraph.text):
                    paragraph.text = paragraph.text.replace(
                        '{{' + template + '}}', ARGS.client[0])

        sdt_logo_stamp.update_logo(slide)

        tags_to_delete = ARGS.tags.split(",")

        if slide.has_notes_slide:
            text_frame = slide.notes_slide.notes_text_frame
            tags_in_comment = sdt_tag_parse.get_all_tags_in_comment(
                text_frame.text)
            matching_tags = copy.copy(tags_in_comment)
            if sdt_common.should_delete_slide(tags_in_comment, tags_to_delete):
                print "Deleting slide {0} with matching tags '{1}'".format(index, matching_tags)
                sdt_common.delete_slide(presentation, slide)
            text_frame.text = sdt_tag_parse.remove_json_from_comment(
                text_frame.text)

    print "List of template strings: {0}".format(templates)
    os.remove(temp_logo_file)
    presentation.save(presentation_file.replace(
        ".ppt", "-{0}.ppt".format(ARGS.client[0])))
    print "Done!"


if __name__ == '__main__':
    main()
